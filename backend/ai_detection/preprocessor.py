"""
ai_detection/preprocessor.py — Input preprocessing pipeline.

Handles: image bytes → numpy array, video → frame sampling,
         text → statistical feature vector, URL → scraped text.
"""

from __future__ import annotations

import asyncio
import io
import logging
from typing import Optional

import numpy as np

logger = logging.getLogger("truthscan.preprocessor")


class ImagePreprocessor:
    """Decode image bytes → RGB numpy array [H, W, 3] uint8."""

    MAX_DIM = 1024  # Downscale anything larger for efficiency

    @staticmethod
    async def from_bytes(data: bytes) -> np.ndarray:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, ImagePreprocessor._decode, data)

    @staticmethod
    def _decode(data: bytes) -> np.ndarray:
        from PIL import Image, ImageOps
        img = Image.open(io.BytesIO(data)).convert("RGB")
        img = ImageOps.exif_transpose(img)  # correct orientation
        # Downscale if needed
        w, h = img.size
        max_dim = ImagePreprocessor.MAX_DIM
        if max(w, h) > max_dim:
            scale = max_dim / max(w, h)
            img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        return np.asarray(img, dtype=np.uint8)


class VideoPreprocessor:
    """Sample representative frames from video for ensemble analysis."""

    N_FRAMES = 8  # analyse 8 evenly-spaced frames

    @staticmethod
    async def extract_frames(data: bytes) -> np.ndarray:
        """
        Extract N_FRAMES evenly-spaced frames and return the
        frame with the most 'information' (highest Laplacian variance).
        """
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, VideoPreprocessor._extract, data)

    @staticmethod
    def _extract(data: bytes) -> np.ndarray:
        try:
            import cv2
            import tempfile, os

            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as f:
                f.write(data)
                tmp_path = f.name

            cap = cv2.VideoCapture(tmp_path)
            total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            if total < 1:
                raise ValueError("Empty video")

            indices = np.linspace(0, total - 1, VideoPreprocessor.N_FRAMES, dtype=int)
            frames = []
            for idx in indices:
                cap.set(cv2.CAP_PROP_POS_FRAMES, int(idx))
                ret, frame = cap.read()
                if ret:
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    frames.append(frame_rgb)
            cap.release()
            os.unlink(tmp_path)

            if not frames:
                raise ValueError("Could not extract frames")

            # Return the sharpest frame (highest Laplacian variance)
            sharpness = [cv2.Laplacian(f, cv2.CV_64F).var() for f in frames]
            return frames[int(np.argmax(sharpness))]

        except ImportError:
            logger.warning("cv2 not available — using PIL for video frame extraction")
            # Fallback: treat first bytes as an image
            return ImagePreprocessor._decode(data[:2_000_000])


class TextPreprocessor:
    """
    For text/file/URL inputs, we extract a 2D feature representation
    that the ensemble models can process. This is a greyscale 'fingerprint'
    image encoding lexical statistics as pixel patterns.
    """

    @staticmethod
    async def extract_from_file(data: bytes, filename: str) -> str:
        """Extract plain text from PDF, DOCX, PPTX, XLSX, TXT."""
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, TextPreprocessor._extract_file, data, filename)

    @staticmethod
    def _extract_file(data: bytes, filename: str) -> str:
        fname = (filename or "").lower()
        try:
            if fname.endswith(".pdf"):
                import pdfplumber
                with pdfplumber.open(io.BytesIO(data)) as pdf:
                    return "\n".join(p.extract_text() or "" for p in pdf.pages[:20])
            if fname.endswith((".docx", ".doc")):
                import docx
                doc = docx.Document(io.BytesIO(data))
                return "\n".join(p.text for p in doc.paragraphs)
            if fname.endswith((".pptx", ".ppt")):
                from pptx import Presentation
                prs = Presentation(io.BytesIO(data))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text_frame.text)
                return "\n".join(texts)
            # Default: try to decode as text
            return data.decode("utf-8", errors="ignore")
        except Exception as e:
            logger.warning("File text extraction failed: %s", e)
            return data.decode("utf-8", errors="ignore")

    @staticmethod
    async def fetch_url(url: str) -> str:
        """Fetch URL and extract visible text."""
        try:
            import httpx
            from bs4 import BeautifulSoup
            async with httpx.AsyncClient(timeout=10, follow_redirects=True) as client:
                resp = await client.get(url, headers={"User-Agent": "TruthScanBot/2.0"})
                soup = BeautifulSoup(resp.text, "html.parser")
                for tag in soup(["script", "style", "nav", "footer"]):
                    tag.decompose()
                return soup.get_text(separator=" ", strip=True)[:20000]
        except Exception as e:
            logger.warning("URL fetch failed: %s", e)
            return f"[URL content unavailable: {e}]"

    @staticmethod
    async def text_to_image(text: str) -> np.ndarray:
        """
        Convert text to a 224×224 statistical 'fingerprint' image.
        Encodes: character n-gram frequencies, word length distribution,
                 punctuation density, sentence length variance.
        This allows image-based models to process text features.
        """
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, TextPreprocessor._encode, text)

    @staticmethod
    def _encode(text: str) -> np.ndarray:
        import re

        words = text.split()
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if s.strip()]

        if not words:
            return np.zeros((224, 224, 3), dtype=np.uint8)

        # Feature extraction
        word_lengths = [len(w) for w in words]
        sent_lengths = [len(s.split()) for s in sentences] if sentences else [0]
        char_counts = {}
        for c in text[:5000]:
            char_counts[c] = char_counts.get(c, 0) + 1

        # Build a 224×224 image encoding statistical features
        img = np.zeros((224, 224, 3), dtype=np.uint8)

        # Channel 0: word length histogram
        for i, length in enumerate(word_lengths[:224]):
            val = min(223, max(0, int((length / 20) * 224)))
            img[i % 224, val, 0] = 200

        # Channel 1: char frequency heatmap
        for i, (char, count) in enumerate(sorted(char_counts.items())[:224]):
            val = min(223, int((count / max(char_counts.values())) * 224))
            img[val, i % 224, 1] = 180

        # Channel 2: sentence length distribution
        for i, sl in enumerate(sent_lengths[:224]):
            val = min(223, max(0, int((sl / 100) * 224)))
            img[i % 224, val, 2] = 160

        return img
